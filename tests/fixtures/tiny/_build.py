"""Synthesize tiny binary fixtures (pdf/docx/pptx/html) for parser tests.

Run once to materialize sample files; the generated binaries are checked
into the repo so tests do not depend on this script at runtime.
"""

from __future__ import annotations

from pathlib import Path

HERE = Path(__file__).resolve().parent


def _tiny_png_bytes(width: int = 320, height: int = 320) -> bytes:
    """Build a small RGB PNG with enough entropy to stay above the 2 KB
    floor after pymupdf re-encodes it into the output PDF."""
    import random
    import struct
    import zlib

    rng = random.Random(1234)
    rows: list[bytes] = []
    for y in range(height):
        row = bytearray([0])
        for x in range(width):
            # Deterministic pseudo-random RGB so each pixel differs; this
            # keeps zlib/flate from collapsing the payload to near nothing.
            r = rng.randrange(256)
            g = rng.randrange(256)
            b = rng.randrange(256)
            row += bytes([r, g, b])
        rows.append(bytes(row))
    raw = b"".join(rows)
    compressed = zlib.compress(raw, 9)

    def _chunk(tag: bytes, data: bytes) -> bytes:
        return (
            struct.pack(">I", len(data))
            + tag
            + data
            + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)
        )

    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    return sig + _chunk(b"IHDR", ihdr) + _chunk(b"IDAT", compressed) + _chunk(b"IEND", b"")


def build_pdf() -> None:
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text(
        (72, 72),
        "Sample PDF\n\nKeywords: photocatalysis, water splitting\n\n"
        "Photocatalysis drives water splitting on TiO2 surfaces.",
        fontsize=11,
    )
    png = _tiny_png_bytes()
    rect = fitz.Rect(72, 200, 272, 400)
    page.insert_image(rect, stream=png)
    page.insert_text(
        (72, 420),
        "Figure 1. Checker pattern used as a synthetic fixture image.",
        fontsize=10,
    )
    doc.save(str(HERE / "sample.pdf"))
    doc.close()


def build_docx() -> None:
    import io

    from docx import Document
    from docx.shared import Inches

    d = Document()
    d.add_heading("Sample DOCX", level=1)
    d.add_paragraph("Keywords: atomic layer deposition, thin films")
    d.add_heading("Introduction", level=2)
    d.add_paragraph("Atomic layer deposition grows thin films one monolayer at a time.")
    d.add_picture(io.BytesIO(_tiny_png_bytes()), width=Inches(2.0))
    d.save(str(HERE / "sample.docx"))


def build_pptx() -> None:
    import io

    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[5]
    s1 = prs.slides.add_slide(blank)
    s1.shapes.title.text = "Sample PPTX"
    tx = s1.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
    tx.text_frame.text = "Photocatalysis converts sunlight into chemical fuel."
    s1.shapes.add_picture(
        io.BytesIO(_tiny_png_bytes()), Inches(4), Inches(4), width=Inches(2), height=Inches(2)
    )
    s2 = prs.slides.add_slide(blank)
    s2.shapes.title.text = "Conclusions"
    tx2 = s2.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
    tx2.text_frame.text = "Atomic layer deposition offers precise thickness control."
    prs.save(str(HERE / "sample.pptx"))


def build_html() -> None:
    html = """<!doctype html>
<html>
<head>
  <title>Sample HTML</title>
  <meta name="author" content="Jane Doe">
  <meta name="date" content="2023-05-01">
  <meta name="description" content="A tiny HTML fixture about semiconductors.">
</head>
<body>
  <h1>Sample HTML</h1>
  <p>Semiconductors are the foundation of modern electronics.
  Photocatalysis on semiconductor surfaces is an active research area.
  This is a short paragraph with enough content for trafilatura to
  extract meaningfully. Doped semiconductors underpin many devices.</p>
  <img src="/images/figure1.png" alt="a figure">
</body>
</html>
"""
    (HERE / "sample.html").write_text(html, encoding="utf-8")


if __name__ == "__main__":
    build_pdf()
    build_docx()
    build_pptx()
    build_html()
    print("fixtures built in", HERE)
