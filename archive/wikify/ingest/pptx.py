"""PPTX ingestion: converts PowerPoint slides to markdown, then reuses extraction pipeline."""

from __future__ import annotations

import hashlib
import json
from pathlib import Path

from pptx import Presentation
from rich.console import Console

from wikify.core.store.models import Paper
from wikify.ingest.extract.chunker import chunk_sections
from wikify.ingest.extract.figure_refs import extract_figure_refs
from wikify.ingest.extract.metadata import _extract_summary, _parse_filename
from wikify.ingest.pdf import ParsedPaper, _parse_section_tree, persist_parsed

console = Console()


# ── Markdown conversion ───────────────────────────────────────────────────────


def _table_to_markdown(table) -> str:
    """Convert a python-pptx Table shape to a simple markdown table."""
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ""

    col_count = max(len(r) for r in rows)
    # Pad rows to equal column count
    rows = [r + [""] * (col_count - len(r)) for r in rows]

    header = "| " + " | ".join(rows[0]) + " |"
    separator = "| " + " | ".join(["---"] * col_count) + " |"
    body_lines = ["| " + " | ".join(row) + " |" for row in rows[1:]]

    return "\n".join([header, separator] + body_lines)


def _shape_text(shape) -> str:
    """Extract plain text from a shape, handling tables specially."""
    if shape.has_table:
        return _table_to_markdown(shape.table)
    if shape.has_text_frame:
        paragraphs = []
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
        return "\n".join(paragraphs)
    return ""


def _is_title_placeholder(shape) -> bool:
    """Return True if the shape is a title-type placeholder."""
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER

        ph = shape.placeholder_format
        if ph is None:
            return False
        return ph.type in (
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
            PP_PLACEHOLDER.SUBTITLE,
        )
    except Exception:
        return False


def _pptx_to_markdown(path: Path) -> str:
    """Convert a PPTX file to markdown text.

    Each slide becomes a ## heading section. Title placeholders set the heading.
    Body text frames follow as paragraphs. Slide notes appear as blockquotes.
    Tables are rendered as markdown tables.
    """
    prs = Presentation(str(path))
    sections: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        # ── Collect title and body shapes ────────────────────────────────────
        title_text: str | None = None
        body_parts: list[str] = []

        for shape in slide.shapes:
            if not (shape.has_text_frame or shape.has_table):
                continue

            if _is_title_placeholder(shape) and shape.has_text_frame:
                candidate = shape.text_frame.text.strip()
                if candidate and title_text is None:
                    title_text = candidate
            else:
                text = _shape_text(shape)
                if text:
                    body_parts.append(text)

        heading = title_text if title_text else f"Slide {slide_idx}"
        slide_label = f"Slide {slide_idx}: {heading}" if title_text else heading

        lines: list[str] = [f"## {slide_label}"]

        if body_parts:
            lines.append("")
            lines.append("\n\n".join(body_parts))

        # ── Slide notes ──────────────────────────────────────────────────────
        try:
            notes_slide = slide.notes_slide
            if notes_slide:
                notes_text = notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    lines.append("")
                    # Prefix each note line with blockquote marker
                    note_lines = "\n".join(
                        f"> {note_line}"
                        for note_line in notes_text.splitlines()
                        if note_line.strip()
                    )
                    lines.append(
                        f"> Note: {note_lines.lstrip('> ')}"
                        if "\n" not in note_lines
                        else note_lines
                    )
        except Exception:
            pass

        sections.append("\n".join(lines))

    return "\n\n".join(sections)


# ── Parsing ───────────────────────────────────────────────────────────────────


def parse_pptx(path: Path) -> ParsedPaper:
    """Parse a PPTX file into structured data. Does NOT touch the database or vault."""
    file_bytes = path.read_bytes()
    file_hash = hashlib.sha256(file_bytes).hexdigest()

    # Convert slides to markdown
    md_text = _pptx_to_markdown(path)

    # Open presentation for metadata
    prs = Presentation(str(path))
    props = prs.core_properties

    # Parse filename for fallback values
    fn_year, fn_author, fn_title = _parse_filename(path.name)

    # ── Title ────────────────────────────────────────────────────────────────
    cp_title = (props.title or "").strip()
    if cp_title:
        title = cp_title
    elif fn_title:
        title = fn_title
    else:
        # Fall back to first slide title extracted from markdown
        first_heading: str | None = None
        for line in md_text.splitlines():
            stripped = line.strip()
            if stripped.startswith("## Slide 1: "):
                first_heading = stripped[len("## Slide 1: ") :].strip()
                break
        title = first_heading or path.stem

    # ── Authors ──────────────────────────────────────────────────────────────
    cp_author = (props.author or "").strip()
    if cp_author:
        authors: list[str] = [
            a.strip() for a in cp_author.replace(";", ",").split(",") if a.strip()
        ]
    elif fn_author:
        authors = [fn_author]
    else:
        authors = []

    # ── Summary ──────────────────────────────────────────────────────────────
    summary = _extract_summary(md_text)

    # ── Year ─────────────────────────────────────────────────────────────────
    year: int | None = fn_year
    if year is None:
        for dt_prop in (props.created, props.modified):
            if dt_prop is not None:
                try:
                    year = dt_prop.year
                    break
                except Exception:
                    pass

    # ── Build Paper record ────────────────────────────────────────────────────
    section_tree = _parse_section_tree(md_text)

    paper = Paper(
        id=file_hash,
        title=title,
        authors=json.dumps(authors),
        summary=summary,
        year=year,
        doi=None,
        doc_type="presentation",
        source_path=str(path),
        file_hash=file_hash,
        section_tree=json.dumps(section_tree),
    )

    # ── Chunks ────────────────────────────────────────────────────────────────
    chunks = chunk_sections(md_text, section_tree, paper.id)

    # ── Figure refs (presentations may mention "Fig. X" in slide text) ────────
    figure_refs = extract_figure_refs(md_text, paper.id)

    return ParsedPaper(
        paper=paper,
        chunks=chunks,
        figures=[],
        citations=[],
        figure_refs=figure_refs,
        md_text=md_text,
    )


# ── Ingestion ─────────────────────────────────────────────────────────────────


def ingest_pptx(path: Path, return_id: bool = False) -> int | str | None:
    """Ingest a single PPTX into the knowledge base.

    Returns 1/0 or the paper ID string if return_id=True (None on skip).
    """
    file_bytes = path.read_bytes()
    file_hash = hashlib.sha256(file_bytes).hexdigest()

    from wikify.core.store.db import get_session

    with get_session() as session:
        existing = session.get(Paper, file_hash)
        if existing:
            console.print(f"[dim]Skipping (unchanged):[/dim] {path.name}")
            return None if return_id else 0

    parsed = parse_pptx(path)
    persist_parsed(parsed)

    console.print(
        f"[green]Ingested:[/green] {path.name} "
        f"({len(parsed.chunks)} chunks, {len(parsed.figure_refs)} figure refs)"
    )
    return parsed.paper.id if return_id else 1
