"""Export slide plans to PowerPoint presentations."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from rich.console import Console

console = Console()

# ── Colour palette ────────────────────────────────────────────────────────────
_NAVY = RGBColor(0x1B, 0x2A, 0x4E)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xC8)
_DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
_TEAL = RGBColor(0x2E, 0x86, 0xAB)

_TEMPLATE_PATH = Path(__file__).parent / "templates" / "professional.pptx"


# ── Backward-compat helper ────────────────────────────────────────────────────


def export_slides(slides: list[dict], output_path: Path, title: str = "") -> Path:
    """Export a list of slide dicts to a PPTX file.

    Each slide dict has: title, bullets, notes, source_papers (optional).
    """
    return PptxExporter().export(slides, output_path, title)


# ── PptxExporter class ────────────────────────────────────────────────────────


class PptxExporter:
    """Export slide plans to a professionally styled PPTX.

    Loads ``src/scholarforge/export/templates/professional.pptx`` when
    available; otherwise applies the same look programmatically so the output
    is never unstyled.

    Slide dict schema (each item in the *slides* list)::

        {
            "title":         str,            # required
            "bullets":       list[str],      # optional
            "notes":         str,            # optional speaker notes
            "source_papers": list[str],      # optional, appended to notes
        }
    """

    def export(
        self,
        slides: list[dict],
        output_path: Path | str,
        title: str = "",
    ) -> Path:
        """Build and save the PPTX file, returning the output path."""
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        if _TEMPLATE_PATH.exists():
            prs = self._export_from_template(slides)
        else:
            console.print("[yellow]Template not found — using programmatic styling.[/yellow]")
            prs = self._export_programmatic(slides)

        prs.save(str(output_path))
        console.print(f"[green]Slides saved:[/green] {output_path}")
        return output_path

    # ── Template-based path ───────────────────────────────────────────────────

    def _export_from_template(self, slides: list[dict]) -> Presentation:
        """Build presentation using the professional.pptx template."""
        prs = Presentation(str(_TEMPLATE_PATH))
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for i, slide_data in enumerate(slides):
            if i == 0:
                self._add_title_slide_template(prs, slide_data)
            else:
                self._add_content_slide_template(prs, slide_data)

        return prs

    def _add_title_slide_template(self, prs: Presentation, data: dict) -> None:
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = data.get("title", "Presentation")
        self._style_run(
            slide.shapes.title.text_frame.paragraphs[0],
            bold=True,
            size=44,
            colour=_WHITE,
        )

        try:
            ph = slide.placeholders[1]
            ph.text = "\n".join(data.get("bullets", []))
            self._style_run(
                ph.text_frame.paragraphs[0],
                bold=False,
                size=24,
                colour=_LIGHT_GRAY,
            )
        except KeyError:
            pass

        if data.get("notes"):
            slide.notes_slide.notes_text_frame.text = data["notes"]

    def _add_content_slide_template(self, prs: Presentation, data: dict) -> None:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = data.get("title", "")
        self._style_run(
            slide.shapes.title.text_frame.paragraphs[0],
            bold=True,
            size=32,
            colour=_WHITE,
        )

        try:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            self._fill_bullets(tf, data.get("bullets", []))
        except KeyError:
            pass

        self._write_notes(slide, data)

    # ── Programmatic fallback ─────────────────────────────────────────────────

    def _export_programmatic(self, slides: list[dict]) -> Presentation:
        """Build presentation with inline styling (no template file needed)."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for i, slide_data in enumerate(slides):
            if i == 0:
                self._add_title_slide_prog(prs, slide_data)
            else:
                self._add_content_slide_prog(prs, slide_data)

        return prs

    def _add_title_slide_prog(self, prs: Presentation, data: dict) -> None:
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        # Navy background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _NAVY

        title_ph = slide.shapes.title
        title_ph.text = data.get("title", "Presentation")
        self._style_run(
            title_ph.text_frame.paragraphs[0],
            bold=True,
            size=44,
            colour=_WHITE,
            align=PP_ALIGN.CENTER,
        )

        try:
            ph = slide.placeholders[1]
            ph.text = "\n".join(data.get("bullets", []))
            self._style_run(
                ph.text_frame.paragraphs[0],
                bold=False,
                size=24,
                colour=_LIGHT_GRAY,
                align=PP_ALIGN.CENTER,
            )
        except KeyError:
            pass

        if data.get("notes"):
            slide.notes_slide.notes_text_frame.text = data["notes"]

    def _add_content_slide_prog(self, prs: Presentation, data: dict) -> None:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        # White background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _WHITE

        # Navy title bar rectangle behind title
        title_bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0),
            Inches(0),
            Inches(13.333),
            Inches(1.2),
        )
        fill = title_bar.fill
        fill.solid()
        fill.fore_color.rgb = _NAVY
        title_bar.line.fill.background()

        title_ph = slide.shapes.title
        title_ph.left = Inches(0.4)
        title_ph.top = Inches(0.1)
        title_ph.width = Inches(12.533)
        title_ph.height = Inches(1.0)
        title_ph.text = data.get("title", "")
        self._style_run(
            title_ph.text_frame.paragraphs[0],
            bold=True,
            size=32,
            colour=_WHITE,
            align=PP_ALIGN.LEFT,
        )

        try:
            body = slide.placeholders[1]
            body.left = Inches(0.5)
            body.top = Inches(1.4)
            body.width = Inches(12.333)
            body.height = Inches(5.7)
            tf = body.text_frame
            tf.clear()
            self._fill_bullets(tf, data.get("bullets", []))
        except KeyError:
            pass

        self._write_notes(slide, data)

    # ── Shared helpers ────────────────────────────────────────────────────────

    @staticmethod
    def _style_run(
        paragraph,
        *,
        bold: bool,
        size: int,
        colour: RGBColor,
        align: PP_ALIGN | None = None,
    ) -> None:
        """Apply font styling to the first run of a paragraph (or add one)."""
        if align is not None:
            paragraph.alignment = align
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        run.font.name = "Calibri"
        run.font.bold = bold
        run.font.size = Pt(size)
        run.font.color.rgb = colour

    @staticmethod
    def _fill_bullets(tf, bullets: list[str]) -> None:
        """Populate a text frame with bullet points."""
        for j, bullet in enumerate(bullets):
            if j == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.text = bullet
            para.space_before = Pt(6)
            run = para.runs[0] if para.runs else para.add_run()
            run.font.name = "Calibri"
            run.font.size = Pt(18)
            run.font.color.rgb = _DARK_GRAY

    @staticmethod
    def _write_notes(slide, data: dict) -> None:
        """Write speaker notes and source papers to the notes pane."""
        parts: list[str] = []
        if data.get("notes"):
            parts.append(data["notes"])
        if data.get("source_papers"):
            parts.append("Sources: " + ", ".join(data["source_papers"]))
        if parts:
            slide.notes_slide.notes_text_frame.text = "\n".join(parts)
