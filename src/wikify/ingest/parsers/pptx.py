"""PPTX parser.

Emits one ``## Slide N`` heading per slide so ``section_spans`` yields
one DocSection per slide. Picture shapes become raw image records in
``metadata['_raw_images']``.
"""

from pathlib import Path

from ..metadata import clean_markdown, extract_summary, parse_filename
from ._sections import section_spans
from .registry import ParseResult


def parse(path: Path) -> ParseResult:
    from pptx import Presentation

    prs = Presentation(str(path))
    md_text = _pptx_to_markdown(prs)
    props = prs.core_properties

    fn_year, fn_author, fn_title = parse_filename(path.name)

    cp_title = (props.title or "").strip()
    if cp_title:
        title = cp_title
    elif fn_title:
        title = fn_title
    else:
        title = path.stem
    title = clean_markdown(title)

    cp_author = (props.author or "").strip()
    if cp_author:
        authors = [a.strip() for a in cp_author.replace(";", ",").split(",") if a.strip()]
    elif fn_author:
        authors = [fn_author]
    else:
        authors = []

    year = fn_year
    if year is None:
        for dt_prop in (props.created, props.modified):
            if dt_prop is not None:
                try:
                    year = dt_prop.year
                    break
                except Exception:
                    pass

    images_raw = _extract_images(prs)
    metadata = {
        "title": title,
        "authors": authors,
        "summary": extract_summary(md_text),
        "year": year,
        "doi": None,
        "_raw_images": images_raw,
    }
    return ParseResult(
        markdown=md_text,
        sections=section_spans(md_text),
        images=[],
        metadata=metadata,
        title=title,
    )


def _is_title_placeholder(shape) -> bool:
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER

        ph = shape.placeholder_format
        if ph is None:
            return False
        return ph.type in (
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
            PP_PLACEHOLDER.SUBTITLE,
        )
    except Exception:
        return False


def _table_to_markdown(table) -> str:
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append(cells)
    if not rows:
        return ""
    col_count = max(len(r) for r in rows)
    rows = [r + [""] * (col_count - len(r)) for r in rows]
    header = "| " + " | ".join(rows[0]) + " |"
    sep = "| " + " | ".join(["---"] * col_count) + " |"
    body = ["| " + " | ".join(r) + " |" for r in rows[1:]]
    return "\n".join([header, sep] + body)


def _shape_text(shape) -> str:
    if shape.has_table:
        return _table_to_markdown(shape.table)
    if shape.has_text_frame:
        paras = []
        for p in shape.text_frame.paragraphs:
            t = p.text.strip()
            if t:
                paras.append(t)
        return "\n".join(paras)
    return ""


def _pptx_to_markdown(prs) -> str:
    sections: list[str] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        title_text: str | None = None
        body_parts: list[str] = []
        for shape in slide.shapes:
            if not (shape.has_text_frame or shape.has_table):
                continue
            if _is_title_placeholder(shape) and shape.has_text_frame:
                cand = shape.text_frame.text.strip()
                if cand and title_text is None:
                    title_text = cand
            else:
                text = _shape_text(shape)
                if text:
                    body_parts.append(text)
        heading = title_text if title_text else f"Slide {slide_idx}"
        slide_label = f"Slide {slide_idx}: {heading}" if title_text else heading
        lines = [f"## {slide_label}"]
        if body_parts:
            lines.append("")
            lines.append("\n\n".join(body_parts))
        try:
            ns = slide.notes_slide
            if ns:
                notes_text = ns.notes_text_frame.text.strip()
                if notes_text:
                    lines.append("")
                    lines.append("> Note: " + notes_text.replace("\n", " "))
        except Exception:
            pass
        sections.append("\n".join(lines))
    return "\n\n".join(sections)


def _extract_images(prs) -> list[dict]:
    raw: list[dict] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_title = ""
        for shape in slide.shapes:
            if _is_title_placeholder(shape) and shape.has_text_frame:
                slide_title = shape.text_frame.text.strip()
                break
        for shape in slide.shapes:
            try:
                if shape.shape_type is None:
                    continue
                if not hasattr(shape, "image"):
                    continue
                img = shape.image
                raw.append(
                    {
                        "bytes": img.blob,
                        "ext": (img.ext or "png").lstrip("."),
                        "page": slide_idx,
                        "caption": slide_title,
                    }
                )
            except Exception:
                continue
    return raw
