"""Create the professional PPTX template for ScholarForge.

Run this script once to generate:
    src/scholarforge/export/templates/professional.pptx

The template defines two slide layouts:
    0 — Title Slide: dark navy background, white title, light-gray subtitle
    1 — Content Slide: white background, navy title bar, dark-gray body text
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x2A, 0x4E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xC8)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
TEAL = RGBColor(0x2E, 0x86, 0xAB)

# ── Dimensions ────────────────────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

OUT_PATH = (
    Path(__file__).parent.parent
    / "src"
    / "scholarforge"
    / "export"
    / "templates"
    / "professional.pptx"
)


def _solid_fill(shape, colour: RGBColor) -> None:
    """Apply a solid fill to a shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _slide_solid_background(slide, colour: RGBColor) -> None:
    """Set a solid background on a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def build_template() -> Path:
    """Build the professional template and return its path."""
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _build_title_slide(prs)
    _build_content_slide(prs)

    prs.save(str(OUT_PATH))
    print(f"Template saved: {OUT_PATH}")
    return OUT_PATH


# ── Layout 0: Title Slide ─────────────────────────────────────────────────────


def _build_title_slide(prs: Presentation) -> None:
    """Build layout-0 Title Slide (dark navy, white title, grey subtitle)."""
    layout = prs.slide_layouts[0]  # "Title Slide" in blank template
    slide = prs.slides.add_slide(layout)

    # Navy background
    _slide_solid_background(slide, NAVY)

    # Title placeholder
    title_ph = slide.shapes.title
    title_ph.left = Inches(1.0)
    title_ph.top = Inches(2.2)
    title_ph.width = Inches(11.333)
    title_ph.height = Inches(1.5)

    tf = title_ph.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.runs[0] if para.runs else para.add_run()
    run.text = "Presentation Title"
    run.font.name = "Calibri"
    run.font.bold = True
    run.font.size = Pt(44)
    run.font.color.rgb = WHITE

    # Subtitle placeholder (idx 1)
    try:
        subtitle_ph = slide.placeholders[1]
        subtitle_ph.left = Inches(1.5)
        subtitle_ph.top = Inches(4.0)
        subtitle_ph.width = Inches(10.333)
        subtitle_ph.height = Inches(1.0)

        stf = subtitle_ph.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sr = sp.runs[0] if sp.runs else sp.add_run()
        sr.text = "Subtitle or Author"
        sr.font.name = "Calibri"
        sr.font.bold = False
        sr.font.size = Pt(24)
        sr.font.color.rgb = LIGHT_GRAY
    except KeyError:
        pass  # subtitle placeholder not present in this layout variant


# ── Layout 1: Content Slide ───────────────────────────────────────────────────


def _build_content_slide(prs: Presentation) -> None:
    """Build layout-1 Content Slide (white background, teal title bar)."""
    layout = prs.slide_layouts[1]  # "Title and Content" in blank template
    slide = prs.slides.add_slide(layout)

    # White background
    _slide_solid_background(slide, WHITE)

    # ── Title bar: teal rectangle behind title text ───────────────────────────
    title_bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0),
        Inches(0),
        SLIDE_WIDTH,
        Inches(1.2),
    )
    _solid_fill(title_bar, NAVY)
    title_bar.line.fill.background()  # no border

    # Title placeholder
    title_ph = slide.shapes.title
    title_ph.left = Inches(0.4)
    title_ph.top = Inches(0.1)
    title_ph.width = Inches(12.533)
    title_ph.height = Inches(1.0)

    ttf = title_ph.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.LEFT
    tr = tp.runs[0] if tp.runs else tp.add_run()
    tr.text = "Slide Title"
    tr.font.name = "Calibri"
    tr.font.bold = True
    tr.font.size = Pt(32)
    tr.font.color.rgb = WHITE

    # Body placeholder (idx 1)
    try:
        body_ph = slide.placeholders[1]
        body_ph.left = Inches(0.5)
        body_ph.top = Inches(1.4)
        body_ph.width = Inches(12.333)
        body_ph.height = Inches(5.7)

        btf = body_ph.text_frame
        btf.word_wrap = True
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.LEFT
        br = bp.runs[0] if bp.runs else bp.add_run()
        br.text = "Body content goes here"
        br.font.name = "Calibri"
        br.font.bold = False
        br.font.size = Pt(20)
        br.font.color.rgb = DARK_GRAY
    except KeyError:
        pass


if __name__ == "__main__":
    build_template()
